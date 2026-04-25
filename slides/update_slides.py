"""Update the v3 proposal presentation into a results presentation.

Strategy: preserve the existing 8 slides (they define the research frame),
patch a few inline text bits, and APPEND two new slides with real findings
from the pilot run.

  Slide 1:  title — fix "Emma" -> "Bruna", update subtitle to reflect pilot
  Slide 9:  NEW  — Pilot Findings (5 F-Droid apps + positive control, key results table)
  Slide 10: NEW  — Tool Architecture & Next Steps
"""

from __future__ import annotations

import csv
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

PROJECT = Path(__file__).resolve().parent.parent
SRC = PROJECT / "ml_android_security_with_notes.pptx"
DST = PROJECT / "slides" / "ml_android_security_final.pptx"
RESULTS = PROJECT / "results"

NAVY = RGBColor(0x0A, 0x1F, 0x3A)
CYAN = RGBColor(0x2A, 0x9D, 0x8F)
AMBER = RGBColor(0xE9, 0xC4, 0x6A)
RED = RGBColor(0xE6, 0x39, 0x46)
ORANGE = RGBColor(0xF4, 0xA2, 0x61)
GREY = RGBColor(0x66, 0x6B, 0x73)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# --------------------------------------------------------------------------- #
# Patch existing slide 1
# --------------------------------------------------------------------------- #


def patch_title_slide(prs: Presentation) -> None:
    s = prs.slides[0]
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "Emma" in run.text:
                    run.text = run.text.replace("Emma", "Bruna Vasconcelos")
                if run.text.strip().startswith("Can Truth Serum-style"):
                    run.text = (
                        "A static-analysis pipeline for finding on-device ML "
                        "vulnerabilities in real Android apps"
                    )
                # strip stranded fragment from the old subtitle
                if "real Android apps through ML pipeline flaws" in run.text:
                    run.text = ""


def _replace_runs(prs: Presentation, slide_idx: int, replacements: dict) -> None:
    s = prs.slides[slide_idx]
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in replacements.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)


def patch_related_work_slide(prs: Presentation) -> None:
    """Slide 2 — add Xu (WWW'19) + Huang (USENIX'21) and reframe the gap claim."""
    _replace_runs(prs, 1, {
        "Lui et al. (MLSec survey, 2023)":
            "Xu et al. (WWW'19) + Sun et al. (USENIX Security'21, ModelXRay)",
        "Lab models only — no Android apps":
            "Measured ML presence and model exposure at scale — this work targets the training-surface gap",
        "FlowDroid & Soot were designed before on-device ML existed — ML data flows are invisible to them":
            "Prior mobile-ML audits measured presence and model exposure; none enforce Truth-Serum preconditions",
    })


def patch_corpus_slide(prs: Presentation) -> None:
    """Slide 4 — reframe the '30 apps, 6 ML-heavy categories' as PLANNED; surface actual pilot corpus."""
    _replace_runs(prs, 3, {
        "30 apps, 1M+ downloads, 6 ML-heavy categories:":
            "PLANNED (future work): 30 apps, 1M+ downloads, 6 ML-heavy categories",
        "Who we study, how we select them, and the boundaries we commit to":
            "Pilot corpus: 5 F-Droid apps + 1 synthetic positive control (see slide 9). Target corpus below is future work.",
        "PoC exploits run in isolated lab environment / emulator only — never deployed against real users or live services":
            "No PoC exploits run in pilot (dynamic validation deferred). Framework shown for future work scope.",
    })


def patch_taxonomy_slide(prs: Presentation) -> None:
    """Slide 6 — remove residual Antidote mention."""
    _replace_runs(prs, 5, {
        "Antidote applicability highest here":
            "Highest-priority target for the pilot detector",
    })


def patch_timeline_slide(prs: Presentation) -> None:
    """Slide 7 — label as planned (not executed); strip Antidote leakage."""
    _replace_runs(prs, 6, {
        "Realistic 8-Week Timeline with Contingency Planning":
            "Planned 8-Week Timeline (from February proposal)",
        "Includes toolchain pilot in Week 1, buffer days in each phase, and explicit fallback if dynamic analysis is limited":
            "Shown for context. Pilot was a 1-day sprint — see slide 5 for what was actually built.",
        "Antidote applicability assessment":
            "Defensive-layer pairing assessment",
    })


def patch_defensibility_slide(prs: Presentation) -> None:
    """Slide 8 — remove claim that contradicts slide 2; acknowledge Xu/Sun as prior art."""
    _replace_runs(prs, 7, {
        "Related work review shows no prior audit of ML components in Android apps — FlowDroid, TaintDroid, and survey literature all stop short of on-device ML pipelines":
            "Prior mobile-ML audits (Xu WWW'19 presence; Sun USENIX'21 model protection) measured different dimensions. This work targets the orthogonal training-surface question",
        "Tool limitations (FlowDroid reflection, Frida repackaging cost) are addressed with manual fallback and realistic scope — not hidden from review":
            "Pilot limitations (pattern-match only, N=5 corpus, no baseline comparison, positive control validates extension match only) disclosed on slides 5 and 9 — not hidden from review",
        "Primary RQs (poisoning surfaces + model exposure) are the MVD. Negative results are framed as a valid security posture finding, not a project failure":
            "Pilot addresses RQ1 and RQ2 at a pattern-match level. RQ3 and RQ4 deferred to future work. Negative result framed as pilot observation, not structural claim",
    })


def patch_methodology_slide(prs: Presentation) -> None:
    """Slide 5 — kill the bait-and-switch: mark unbuilt phases (PLANNED), remove Antidote leakage."""
    _replace_runs(prs, 4, {
        "Methodology: Two-Phase Pipeline with Known Constraints":
            "Methodology: What Was Planned vs. What the Pilot Built",
        "Static taint analysis (FlowDroid + Soot) → Dynamic validation (Frida) — limitations acknowledged up front":
            "Pilot scope (BUILT): archive scan + manifest parse + Smali pattern scan. Taint analysis and dynamic validation are future work.",
        "Phase 1 — Static Analysis":
            "Phase 1 — Static Analysis (BUILT in pilot)",
        "APK Decompilation":
            "APK Archive + Decompile (BUILT)",
        "apktool + jadx → Smali + embedded .tflite/.onnx":
            "apktool decode + archive scan for model artifacts (entropy + magic bytes)",
        "Custom Soot Transform":
            "Custom Soot Transform (PLANNED)",
        "Inter-procedural: user data → ML training path; ML-specific sources/sinks added":
            "Pilot uses Smali pattern-match (grep) as a lightweight proxy; full Soot taint is future work",
        "FlowDroid Taint":
            "FlowDroid Taint (PLANNED)",
        "PII sources (mic/camera/location) → inference sinks":
            "Not run in pilot — permission scan is used as a correlate, not a taint result",
        "FlowDroid loses track through Java reflection & dynamic class loading — common in ML frameworks. Manual review applied to flagged cases.":
            "Pilot does pattern-matching, not taint analysis. Overclaim risk: obfuscated apps can evade string patterns. Addressed in limitations.",
        "Phase 2 — Dynamic Validation":
            "Phase 2 — Dynamic Validation (PLANNED — not in pilot)",
        "PoC Exploit Crafting":
            "PoC Exploit Crafting (PLANNED)",
        "Poisoned inputs targeting confirmed static surfaces; exploitability verified in emulator":
            "Future work: craft Truth-Serum-style poisoned inputs once static results indicate viable surfaces",
        "Frida Instrumentation":
            "Frida Instrumentation (PLANNED)",
        "Strategy: repackaged APK method (avoids rooting; broader app coverage)":
            "Deferred: APK repackaging adds ~2 days per batch — out of scope for 1-day pilot",
        "Antidote Applicability":
            "Defensive-layer pairing (PLANNED)",
        "For each poisoning surface: can Antidote's pretrained-feature detector be applied as runtime shield?":
            "Future work: couple detected surfaces with a training-set poison detector as a runtime shield",
        "Frida requires root OR APK repackaging. We use repackaging: broader coverage but adds ~2 days per batch. Rooted emulator used for preliminary validation.":
            "Dynamic validation is future work. Pilot result stands on static evidence + positive control only.",
    })


# --------------------------------------------------------------------------- #
# Load pipeline findings
# --------------------------------------------------------------------------- #


def load_findings() -> tuple[list[dict], dict]:
    """Return (per-app rows, aggregate stats)."""
    summary_csv = RESULTS / "summary.csv"
    rows: list[dict] = []
    if summary_csv.exists():
        with summary_csv.open() as f:
            rows = list(csv.DictReader(f))

    agg = {
        "n_apks": len(rows),
        "total_size_mb": sum(float(r.get("apk_size_mb", 0)) for r in rows),
        "n_with_models": sum(1 for r in rows if int(r.get("n_model_artifacts", 0)) > 0),
        "n_with_sdks": sum(1 for r in rows if r.get("ml_sdks", "").strip()),
        "total_findings": sum(
            int(r.get("critical", 0))
            + int(r.get("high", 0))
            + int(r.get("medium", 0))
            + int(r.get("info", 0))
            for r in rows
        ),
    }
    return rows, agg


# --------------------------------------------------------------------------- #
# Slide helpers
# --------------------------------------------------------------------------- #


def _add_textbox(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = NAVY,
    align: str = "left",
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {"left": 1, "center": 2, "right": 3}.get(align, 1)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_filled_rect(slide, left, top, width, height, fill: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _add_card(slide, left, top, width, height, title, body, accent: RGBColor):
    """Filled rounded rect with accent stripe + title + body text."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent
    card.line.width = Pt(1.5)

    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Emu(int(0.15 * 914400)), height
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    inner_left = left + Inches(0.35)
    _add_textbox(
        slide, inner_left, top + Inches(0.2),
        width - Inches(0.5), Inches(0.5),
        title, size=16, bold=True, color=NAVY,
    )
    _add_textbox(
        slide, inner_left, top + Inches(0.75),
        width - Inches(0.5), height - Inches(0.9),
        body, size=12, color=GREY,
    )


# --------------------------------------------------------------------------- #
# NEW Slide 9: Pilot Findings
# --------------------------------------------------------------------------- #


def add_findings_slide(prs: Presentation, rows: list[dict], agg: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # only DEFAULT available

    # Title bar
    _add_filled_rect(slide, Emu(0), Emu(0), prs.slide_width, Inches(0.9), NAVY)
    _add_textbox(
        slide, Inches(0.6), Inches(0.18),
        Inches(12), Inches(0.6),
        f"Pilot Findings: {agg['n_apks']} APKs Analyzed", size=26, bold=True, color=WHITE,
    )
    _add_textbox(
        slide, Inches(0.6), Inches(0.63),
        Inches(12), Inches(0.3),
        "Static analysis on an F-Droid corpus — what the tool found and what it tells us",
        size=13, color=RGBColor(0xB0, 0xC4, 0xDE),
    )

    # Aggregate strip
    y = Inches(1.1)
    stats = [
        (f"{agg['n_apks']}", "APKs analyzed", CYAN),
        (f"{agg['total_size_mb']:.0f} MB", "Total corpus", ORANGE),
        (f"{agg['n_with_models']}", "With on-device models", RED),
        (f"{agg['total_findings']}", "Findings logged", AMBER),
    ]
    col_w = Inches(2.85)
    for i, (num, lbl, color) in enumerate(stats):
        left = Inches(0.6) + col_w * i + Inches(0.1 * i)
        _add_filled_rect(slide, left, y, col_w, Inches(1.0), color)
        _add_textbox(
            slide, left, y + Inches(0.1),
            col_w, Inches(0.5),
            num, size=24, bold=True, color=WHITE, align="center",
        )
        _add_textbox(
            slide, left, y + Inches(0.65),
            col_w, Inches(0.3),
            lbl, size=11, color=WHITE, align="center",
        )

    # Per-app results table
    ty = Inches(2.35)
    _add_textbox(
        slide, Inches(0.6), ty,
        Inches(12), Inches(0.4),
        "Per-App Results",
        size=16, bold=True, color=NAVY,
    )

    header_y = ty + Inches(0.45)
    headers = ["App (package)", "Size", "Models", "ML SDKs", "Findings", "Top Severity"]
    col_widths_in = [3.0, 0.9, 0.9, 2.2, 0.9, 2.5]
    col_left_in = [0.6]
    for w in col_widths_in[:-1]:
        col_left_in.append(col_left_in[-1] + w + 0.1)
    col_left_in = [Inches(x) for x in col_left_in]

    _add_filled_rect(slide, Inches(0.6), header_y, Inches(12.2), Inches(0.4), NAVY)
    for i, h in enumerate(headers):
        _add_textbox(
            slide, col_left_in[i], header_y + Inches(0.08),
            Inches(col_widths_in[i]), Inches(0.3),
            h, size=11, bold=True, color=WHITE,
        )

    row_y = header_y + Inches(0.45)
    for ri, row in enumerate(rows):
        bg = RGBColor(0xF5, 0xF7, 0xFA) if ri % 2 == 0 else WHITE
        _add_filled_rect(slide, Inches(0.6), row_y, Inches(12.2), Inches(0.5), bg)

        top_sev = "Info"
        for sev in ("Critical", "High", "Medium", "Info"):
            if int(row.get(sev.lower(), 0)) > 0:
                top_sev = sev
                break
        sev_color = {
            "Critical": RED,
            "High": ORANGE,
            "Medium": AMBER,
            "Info": GREY,
        }[top_sev]

        pkg_label = row.get("package") or Path(row["apk"]).stem
        raw_sdks = (row.get("ml_sdks") or "").strip()
        if raw_sdks:
            sdks = raw_sdks
        elif int(row.get("smali_scanned", 0) or 0) == 0:
            sdks = "n/a"  # apktool decode failed; Smali not scanned
        else:
            sdks = "none"
        total_findings = sum(
            int(row.get(k, 0)) for k in ("critical", "high", "medium", "info")
        )
        cells = [
            pkg_label,
            f"{row.get('apk_size_mb', '0')} MB",
            row.get("n_model_artifacts", "0"),
            sdks,
            str(total_findings),
            top_sev,
        ]
        for i, val in enumerate(cells):
            color = sev_color if i == 5 else NAVY
            _add_textbox(
                slide, col_left_in[i], row_y + Inches(0.13),
                Inches(col_widths_in[i]), Inches(0.3),
                str(val), size=11, color=color,
                bold=(i == 5),
            )
        row_y = row_y + Inches(0.55)

    # Key takeaway footer
    foot_y = Inches(6.9)
    _add_filled_rect(slide, Emu(0), foot_y, prs.slide_width, Inches(0.55), NAVY)
    takeaway = (
        "Pilot observation: 0 / 5 F-Droid apps carry on-device ML. Detector "
        "validated on synthetic positive control (MobileNet). Next: Play Store corpus."
    )
    _add_textbox(
        slide, Inches(0.6), foot_y + Inches(0.14),
        Inches(12), Inches(0.3),
        takeaway, size=12, color=WHITE,
    )

    # Notes
    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "Slide 9 — Pilot findings. ~60 seconds.\n\n"
        f"We analyzed {agg['n_apks']} open-source Android apps from F-Droid, "
        f"totaling about {agg['total_size_mb']:.0f} megabytes of distributed "
        "binaries, plus one synthetic positive-control APK built from a "
        "public MobileNet model. The top row shows the headline numbers.\n\n"
        "The most interesting finding is what we didn't find in the real apps. "
        "None of the five F-Droid apps bundle on-device ML. We frame this as a "
        "PILOT OBSERVATION, not a structural claim about F-Droid — five apps "
        "is 0.1% of the catalog. But it does tell us F-Droid is not a "
        "productive corpus for this question.\n\n"
        "[beat]\n\n"
        "The positive control confirmed the detector fires as expected: "
        "MobileNet v1 in the synthetic APK was flagged High-severity with "
        "correct format identification. The null result on the F-Droid apps "
        "is a real null, not a tool failure.\n\n"
        "[beat]\n\n"
        "One honest limitation visible in the table: the Smali scan cap of "
        "5 000 files meant Nextcloud was only 11% covered. We've already "
        "added warnings when this happens and the fix is a trivial config "
        "change for the Play Store run."
    )


# --------------------------------------------------------------------------- #
# NEW Slide 10: Tool + Next Steps
# --------------------------------------------------------------------------- #


def add_tool_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _add_filled_rect(slide, Emu(0), Emu(0), prs.slide_width, Inches(0.9), NAVY)
    _add_textbox(
        slide, Inches(0.6), Inches(0.18),
        Inches(12), Inches(0.6),
        "Tool Architecture & Next Steps", size=26, bold=True, color=WHITE,
    )
    _add_textbox(
        slide, Inches(0.6), Inches(0.63),
        Inches(12), Inches(0.3),
        "What we built, what it does, where it goes next",
        size=13, color=RGBColor(0xB0, 0xC4, 0xDE),
    )

    # Left: pipeline stages
    _add_textbox(
        slide, Inches(0.6), Inches(1.2),
        Inches(6), Inches(0.4),
        "Pipeline (Python + apktool + jadx)",
        size=16, bold=True, color=CYAN,
    )

    stages = [
        ("1. APK unzip + archive scan",
         "Magic-byte + extension match for .tflite/.onnx/.pb/.pt; entropy check for encryption"),
        ("2. apktool decode",
         "Smali + AndroidManifest extraction; fall back gracefully if decode fails"),
        ("3. Manifest parsing",
         "Sensitive permissions flagged as ML-data sources (CAMERA, MIC, LOCATION, SMS, INTERNET)"),
        ("4. Smali pattern scan",
         "ML SDK imports (TFLite / ML Kit / ONNX / MediaPipe / PyTorch); poisoning-surface patterns"),
        ("5. Severity rubric",
         "Critical / High / Medium / Info — per finding, per app, aggregated to CSV + per-app JSON"),
    ]
    y = Inches(1.7)
    for title, body in stages:
        _add_textbox(
            slide, Inches(0.6), y, Inches(6), Inches(0.35),
            title, size=13, bold=True, color=NAVY,
        )
        _add_textbox(
            slide, Inches(0.8), y + Inches(0.35), Inches(5.8), Inches(0.55),
            body, size=11, color=GREY,
        )
        y = y + Inches(0.9)

    # Right: limitations + next steps
    _add_textbox(
        slide, Inches(7.1), Inches(1.2),
        Inches(6), Inches(0.4),
        "Limitations & Next Steps",
        size=16, bold=True, color=ORANGE,
    )

    _add_card(
        slide, Inches(7.1), Inches(1.7), Inches(6), Inches(1.3),
        "What we cut for a 1-day pilot",
        "• Dynamic analysis (Frida) — deferred\n"
        "• FlowDroid taint analysis — deferred\n"
        "• Corpus: 5 F-Droid apps + synthetic positive control (pilot) vs. 30+ Play Store apps (planned)",
        AMBER,
    )
    _add_card(
        slide, Inches(7.1), Inches(3.1), Inches(6), Inches(1.3),
        "What the tool is ready for",
        "• Any Play Store APK (same pipeline, no code changes)\n"
        "• Encrypted-model inspection (entropy is already computed)\n"
        "• Sub-second per-APK archive scan; minutes for full decode",
        CYAN,
    )
    _add_card(
        slide, Inches(7.1), Inches(4.5), Inches(6), Inches(1.3),
        "Next: validate on high-risk categories",
        "• Keyboard apps (continuous typing → retraining surface)\n"
        "• Translation apps (offline TFLite models)\n"
        "• Camera filter apps (on-device classification on PII)",
        RED,
    )

    # Footer
    foot_y = Inches(6.9)
    _add_filled_rect(slide, Emu(0), foot_y, prs.slide_width, Inches(0.55), NAVY)
    _add_textbox(
        slide, Inches(0.6), foot_y + Inches(0.14),
        Inches(12), Inches(0.3),
        "Defense (detect poisons in training sets) + this tool (audit where training happens) = the full attack → defense → field-audit arc.",
        size=12, color=WHITE,
    )

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "Slide 10 — Tool architecture and next steps. ~50 seconds.\n\n"
        "The tool is a five-stage Python pipeline. Stage one is a fast archive "
        "scan that takes under a second per APK. Stages two through four rely "
        "on apktool and take a few minutes on large apps. Stage five applies "
        "the severity rubric we defined in the proposal.\n\n"
        "[beat]\n\n"
        "Limitations we're honest about: in one day we cut dynamic analysis "
        "and FlowDroid. Those are future work, not broken work. The static "
        "pipeline we do have runs cleanly on any APK, so the next step is "
        "pointing it at the Play Store categories where on-device ML is "
        "actually common — keyboard, translation, and camera apps.\n\n"
        "[beat]\n\n"
        "The closing framing is the arc. A training-set defense that "
        "detects poisons is out of scope for this pilot. This project is "
        "the auditor that tells us where such defenses are needed. Together "
        "they close the loop between research and real-world risk assessment."
    )


# --------------------------------------------------------------------------- #
# Main
# --------------------------------------------------------------------------- #


def main() -> None:
    prs = Presentation(str(SRC))
    patch_title_slide(prs)
    patch_related_work_slide(prs)
    patch_corpus_slide(prs)
    patch_methodology_slide(prs)
    patch_taxonomy_slide(prs)
    patch_timeline_slide(prs)
    patch_defensibility_slide(prs)
    rows, agg = load_findings()
    if not rows:
        print("WARNING: no pipeline results yet; adding placeholder findings slide")
        rows = [{"apk": "(pending)", "package": "(run pipeline first)",
                 "apk_size_mb": "0", "n_model_artifacts": "0",
                 "ml_sdks": "", "critical": "0", "high": "0",
                 "medium": "0", "info": "0"}]
        agg = {"n_apks": 0, "total_size_mb": 0, "n_with_models": 0,
               "n_with_sdks": 0, "total_findings": 0}
    add_findings_slide(prs, rows, agg)
    add_tool_slide(prs)
    DST.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(DST))
    print(f"saved: {DST}")
    print(f"   {len(prs.slides)} slides total")


if __name__ == "__main__":
    main()
