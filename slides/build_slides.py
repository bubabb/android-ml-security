"""Build the CS 6491 final presentation from scratch (research-talk structure).

Motivation → Related Work → Threat Model → Approach → Methods (Data, Eval) →
Results (Findings, Coverage chart) → Limitations → Conclusion.

Reads results/summary.csv. Writes slides/ml_android_security_final.pptx.
"""
from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

PROJECT = Path(__file__).resolve().parent.parent
RESULTS = PROJECT / "results"
DST = PROJECT / "slides" / "ml_android_security_final.pptx"

# Palette
NAVY = RGBColor(0x0A, 0x1F, 0x3A)
DARK = RGBColor(0x15, 0x25, 0x3F)
CYAN = RGBColor(0x2A, 0x9D, 0x8F)
TEAL = RGBColor(0x26, 0x7A, 0x6F)
AMBER = RGBColor(0xE9, 0xC4, 0x6A)
ORANGE = RGBColor(0xF4, 0xA2, 0x61)
RED = RGBColor(0xE6, 0x39, 0x46)
GREEN = RGBColor(0x4A, 0x8C, 0x4A)
GREY = RGBColor(0x55, 0x5B, 0x66)
LIGHT_GREY = RGBColor(0xE8, 0xE8, 0xEC)
PALE = RGBColor(0xF6, 0xF7, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xB0, 0xC4, 0xDE)

TOTAL_SLIDES = 11


# ----------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, left, top, width, height, fill, no_line=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if no_line:
        shape.line.fill.background()
    return shape


def rrect(slide, left, top, width, height, fill, no_line=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if no_line:
        shape.line.fill.background()
    return shape


def tbox(slide, left, top, width, height, text, *,
         size=14, bold=False, italic=False,
         color=NAVY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def multi_tbox(slide, left, top, width, height, lines, *,
               size=12, color=NAVY, bullet=True):
    """lines = list of strings (or (bold_prefix, rest) tuples)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        prefix = "•  " if bullet else ""
        if isinstance(line, tuple):
            bold_part, rest = line
            r1 = p.add_run()
            r1.text = prefix + bold_part
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = " " + rest
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = prefix + line
            r.font.size = Pt(size)
            r.font.color.rgb = color
        p.space_after = Pt(4)


def header(slide, prs, title, subtitle=None):
    rect(slide, 0, 0, prs.slide_width, Inches(0.9), NAVY)
    rect(slide, 0, Inches(0.9), prs.slide_width, Inches(0.05), CYAN)
    tbox(slide, Inches(0.6), Inches(0.2), prs.slide_width - Inches(1.2), Inches(0.45),
         title, size=24, bold=True, color=WHITE)
    if subtitle:
        tbox(slide, Inches(0.6), Inches(0.58), prs.slide_width - Inches(1.2), Inches(0.32),
             subtitle, size=12, color=LIGHT_BLUE, italic=True)


def page_num(slide, prs, n):
    tbox(slide, prs.slide_width - Inches(1.1), prs.slide_height - Inches(0.35),
         Inches(0.9), Inches(0.25),
         f"{n} / {TOTAL_SLIDES}", size=9, color=GREY, align=PP_ALIGN.RIGHT)
    tbox(slide, Inches(0.3), prs.slide_width and prs.slide_height - Inches(0.35),
         Inches(10), Inches(0.25),
         "Bruna Vasconcelos  •  CS 6491 Final Project  •  Spring 2026",
         size=9, color=GREY)


def card(slide, left, top, width, height, title, body_lines, *,
         accent=CYAN, title_size=13, body_size=10):
    # Card body
    rrect(slide, left, top, width, height, WHITE)
    # Border (outlined rrect)
    border = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    border.fill.background()
    border.line.color.rgb = LIGHT_GREY
    border.line.width = Pt(1)
    # Accent stripe
    rect(slide, left, top, width, Inches(0.08), accent)
    # Title
    tbox(slide, left + Inches(0.18), top + Inches(0.16),
         width - Inches(0.36), Inches(0.35),
         title, size=title_size, bold=True, color=NAVY)
    # Body
    multi_tbox(
        slide, left + Inches(0.18), top + Inches(0.55),
        width - Inches(0.36), height - Inches(0.65),
        body_lines, size=body_size, color=GREY, bullet=True,
    )


def kpi(slide, left, top, width, height, number, label, color):
    rrect(slide, left, top, width, height, color)
    tbox(slide, left, top + Inches(0.1), width, Inches(0.7),
         number, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tbox(slide, left, top + Inches(0.85), width, Inches(0.3),
         label, size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ----------------------------------------------------------------------
# Load pipeline findings
# ----------------------------------------------------------------------

def load_rows():
    with (RESULTS / "summary.csv").open() as f:
        return list(csv.DictReader(f))


def aggregate(rows):
    return {
        "n_apks": len(rows),
        "real_apps": sum(1 for r in rows if "positive" not in r["apk"].lower()),
        "total_mb": sum(float(r.get("apk_size_mb", 0)) for r in rows),
        "real_mb": sum(
            float(r.get("apk_size_mb", 0)) for r in rows
            if "positive" not in r["apk"].lower()
        ),
        "n_with_models": sum(1 for r in rows if int(r.get("n_model_artifacts", 0)) > 0),
        "total_findings": sum(
            int(r.get("critical", 0)) + int(r.get("high", 0))
            + int(r.get("medium", 0)) + int(r.get("info", 0)) for r in rows
        ),
    }


# ----------------------------------------------------------------------
# Slides
# ----------------------------------------------------------------------

def slide_1_title(prs):
    s = blank(prs)
    # Full-bleed dark background
    rect(s, 0, 0, prs.slide_width, prs.slide_height, DARK)
    # Cyan accent stripe
    rect(s, 0, Inches(3.5), prs.slide_width, Inches(0.06), CYAN)
    # Small course tag
    tbox(s, Inches(1), Inches(2.3), Inches(11.33), Inches(0.3),
         "CS 6491 · SOFTWARE & SYSTEMS SECURITY · FINAL PROJECT",
         size=12, color=CYAN, bold=True)
    # Title
    tbox(s, Inches(1), Inches(2.7), Inches(11.33), Inches(0.8),
         "Security Analysis of ML-Powered Android Apps",
         size=36, bold=True, color=WHITE)
    # Subtitle
    tbox(s, Inches(1), Inches(3.7), Inches(11.33), Inches(0.5),
         "A static-analysis pipeline for auditing the Truth-Serum training-surface precondition",
         size=16, color=LIGHT_BLUE, italic=True)
    # Author line
    tbox(s, Inches(1), Inches(6.2), Inches(11.33), Inches(0.35),
         "Bruna Vasconcelos", size=16, bold=True, color=WHITE)
    tbox(s, Inches(1), Inches(6.55), Inches(11.33), Inches(0.3),
         "University of Utah   ·   Spring 2026",
         size=12, color=LIGHT_BLUE)


def slide_2_motivation(prs):
    s = blank(prs)
    header(s, prs, "Motivation & Problem",
           "Truth Serum (Tramèr et al., CCS 2022) — poisoning that amplifies membership inference. Precondition maps onto mobile ML.")

    # Left column: the attack
    tbox(s, Inches(0.6), Inches(1.1), Inches(6), Inches(0.35),
         "THE TRUTH SERUM ATTACK", size=11, bold=True, color=CYAN)

    steps = [
        ("1", "Attacker contributes\npoisoned training samples", CYAN),
        ("2", "Model trains on\npoisoned + real data", TEAL),
        ("3", "Targets show\nelevated training loss", AMBER),
        ("4", "Adversary queries →\nextracts real-data secrets", RED),
    ]
    step_w = Inches(1.4)
    step_h = Inches(1.2)
    gap = Inches(0.08)
    y = Inches(1.55)
    for i, (n, text, color) in enumerate(steps):
        x = Inches(0.6) + (step_w + gap) * i
        rrect(s, x, y, step_w, step_h, color)
        tbox(s, x, y + Inches(0.08), step_w, Inches(0.35),
             n, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tbox(s, x + Inches(0.08), y + Inches(0.5), step_w - Inches(0.16), step_h - Inches(0.55),
             text, size=10, color=WHITE, align=PP_ALIGN.CENTER)

    tbox(s, Inches(0.6), Inches(2.9), Inches(6.3), Inches(0.4),
         "Precondition: attacker can contribute training data. No root, no network intercept, no code exec needed.",
         size=11, color=GREY, italic=True)

    # Right column: why mobile
    tbox(s, Inches(7.2), Inches(1.1), Inches(5.6), Inches(0.35),
         "ON-DEVICE ML MAKES THIS EASY", size=11, bold=True, color=CYAN)

    items = [
        ("Keyboards", "User input CAN feed into on-device personalisation or federated learning (e.g., GBoard)."),
        ("Translation", "Offline TFLite models; user corrections CAN update local weights (e.g., Pixel on-device translate)."),
        ("Camera filters", "On-device classification on PII; adaptation loops are possible — which is the audit question."),
    ]
    y = Inches(1.55)
    for title, desc in items:
        rrect(s, Inches(7.2), y, Inches(5.6), Inches(0.43), LIGHT_GREY)
        tbox(s, Inches(7.35), y + Inches(0.05), Inches(1.8), Inches(0.33),
             title, size=12, bold=True, color=NAVY)
        tbox(s, Inches(9.0), y + Inches(0.08), Inches(3.7), Inches(0.3),
             desc, size=9, color=GREY)
        y += Inches(0.5)

    # Research question banner
    rrect(s, Inches(0.6), Inches(5.2), prs.slide_width - Inches(1.2), Inches(1.5), NAVY)
    tbox(s, Inches(0.85), Inches(5.3), Inches(12), Inches(0.3),
         "RESEARCH QUESTION", size=11, bold=True, color=CYAN)
    tbox(s, Inches(0.85), Inches(5.6), Inches(12), Inches(0.45),
         "Does the Truth-Serum training-surface precondition exist in real Android apps?",
         size=17, bold=True, color=WHITE)
    tbox(s, Inches(0.85), Inches(6.1), Inches(12), Inches(0.5),
         "And can we detect it with static analysis alone — no rooting, no instrumentation, no user interaction?",
         size=12, color=LIGHT_BLUE, italic=True)

    page_num(s, prs, 2)


def slide_3_related_work(prs):
    s = blank(prs)
    header(s, prs, "Related Work & Gap",
           "Two dimensions already measured at scale — training surface is the one left")

    card_w = Inches(3.02)
    card_h = Inches(2.4)
    y1 = Inches(1.25)
    gap = Inches(0.12)

    # Four cards in a row
    card(s, Inches(0.6) + (card_w + gap) * 0, y1, card_w, card_h,
         "FlowDroid (PLDI '14)",
         [("Scope:", "Android taint analysis — user data to network sinks."),
          ("Gap:", "No ML-specific sources or sinks. TFLite is invisible."),
          ("Relation:", "Predates on-device ML entirely.")],
         accent=GREY)
    card(s, Inches(0.6) + (card_w + gap) * 1, y1, card_w, card_h,
         "TaintDroid (OSDI '10)",
         [("Scope:", "OS-level dynamic taint tracking for PII."),
          ("Gap:", "Older still. No awareness of ML inference or training."),
          ("Relation:", "Complementary; different granularity.")],
         accent=GREY)
    card(s, Inches(0.6) + (card_w + gap) * 2, y1, card_w, card_h,
         "Xu et al. (WWW '19)",
         [("Scope:", "First empirical study — 16,500 Play Store apps."),
          ("Measures:", "Presence of DL apps, frameworks used."),
          ("Gap:", "Presence ≠ training-surface exposure.")],
         accent=TEAL)
    card(s, Inches(0.6) + (card_w + gap) * 3, y1, card_w, card_h,
         "Sun et al. (USENIX '21) — ModelXRay",
         [("Scope:", "Large-scale model-protection study."),
          ("Finding:", "41% unprotected; 66% of protected still extractable."),
          ("Gap:", "Extraction ≠ training-surface.")],
         accent=TEAL)

    # Gap banner
    rrect(s, Inches(0.6), Inches(4.0), prs.slide_width - Inches(1.2), Inches(2.6), NAVY)
    tbox(s, Inches(0.85), Inches(4.15), Inches(12), Inches(0.3),
         "THE GAP THIS WORK TARGETS", size=11, bold=True, color=CYAN)
    tbox(s, Inches(0.85), Inches(4.45), Inches(12), Inches(0.5),
         "None of these tools enforces the Truth-Serum precondition — whether user input flows into on-device training.",
         size=15, bold=True, color=WHITE)
    # Three-column matrix
    labels = [
        ("Presence of ML", "Xu '19  ✓", LIGHT_BLUE),
        ("Model protection", "Sun '21  ✓", LIGHT_BLUE),
        ("Training surface", "This work", AMBER),
    ]
    col_w = Inches(4.0)
    y = Inches(5.1)
    for i, (dim, who, color) in enumerate(labels):
        x = Inches(0.85) + (col_w + Inches(0.1)) * i
        tbox(s, x, y, col_w, Inches(0.35),
             dim, size=12, bold=True, color=WHITE)
        tbox(s, x, y + Inches(0.35), col_w, Inches(0.35),
             who, size=13, bold=True, color=color)
    tbox(s, Inches(0.85), Inches(6.2), Inches(12), Inches(0.3),
         "This work complements Xu and Sun rather than competing with them — orthogonal audit dimension.",
         size=11, italic=True, color=LIGHT_BLUE)

    page_num(s, prs, 3)


def slide_4_threat_model(prs):
    s = blank(prs)
    header(s, prs, "Threat Model",
           "Precise adversary capabilities, defender constraints, and what is out of scope")

    col_w = Inches(3.02)
    col_h = Inches(5.5)
    y0 = Inches(1.25)
    gap = Inches(0.12)

    card(s, Inches(0.6) + (col_w + gap) * 0, y0, col_w, col_h,
         "Adversary",
         [("Role:", "Malicious training-data contributor"),
          ("Capability:", "Influences data ingested by on-device training, personalisation, or federated-learning loop"),
          ("Does NOT need:", "Code execution on the device, root, or network-level interception"),
          ("Example:", "A Truth-Serum-style poisoned input delivered through the normal data-collection channel")],
         accent=RED)
    card(s, Inches(0.6) + (col_w + gap) * 1, y0, col_w, col_h,
         "Defender",
         [("Role:", "App auditor"),
          ("Access:", "Publicly distributed APK only"),
          ("Constraint:", "No running instance, no backend, no runtime instrumentation"),
          ("Method:", "Static analysis only (matches ModelXRay and FlowDroid setting)")],
         accent=CYAN)
    card(s, Inches(0.6) + (col_w + gap) * 2, y0, col_w, col_h,
         "Assets at Risk",
         [("Model weights:", "Integrity (poisoning) and confidentiality (extraction)"),
          ("Training pipeline:", "Integrity — no adversarial input filtering"),
          ("User PII:", "Routed through inference / training loops"),
          ("Severity driver:", "Presence of the training surface raises all three")],
         accent=AMBER)
    card(s, Inches(0.6) + (col_w + gap) * 3, y0, col_w, col_h,
         "Out of Scope",
         [("Runtime exploits:", "Memory-corruption in native TFLite libraries"),
          ("Supply chain:", "Compromised model dependencies"),
          ("Backend services:", "Server-side training pipeline attacks"),
          ("Repackaging:", "Device-side model replacement (covered by ModelXRay)")],
         accent=GREY)

    page_num(s, prs, 4)


def slide_5_approach(prs):
    s = blank(prs)
    header(s, prs, "Approach: Five-Stage Static Pipeline",
           "Architected to run on any APK without running the app — outputs per-app JSON + aggregated CSV")

    # Pipeline stages as horizontal flow
    stages = [
        ("1", "Archive Scan",
         "unzip APK · extension match\n(.tflite/.onnx/.pb/.pt)\nShannon entropy · magic bytes",
         CYAN),
        ("2", "Native Libs",
         "libtensorflowlite.so\nlibonnxruntime.so\nlibmediapipe.so / etc.",
         TEAL),
        ("3", "apktool Decode",
         "Smali bytecode\nAndroidManifest.xml\n300s timeout · safe_rmtree",
         AMBER),
        ("4", "Manifest Parse",
         "defusedxml (XXE-safe)\nCAMERA · MIC · LOCATION\nINTERNET · SMS · STORAGE",
         ORANGE),
        ("5", "Smali Patterns",
         "ML-SDK imports\nPoisoning-surface regexes\nCap: 5,000 files/app",
         RED),
    ]
    stage_w = Inches(2.35)
    stage_h = Inches(2.1)
    y = Inches(1.3)
    total_w = stage_w.emu * 5 + Emu(Inches(0.15).emu * 4).emu
    start_x = (prs.slide_width - Emu(total_w)) / 2

    for i, (n, title, body, color) in enumerate(stages):
        x = start_x + (stage_w.emu + Inches(0.15).emu) * i
        rrect(s, Emu(x), y, stage_w, stage_h, color)
        # Number circle
        tbox(s, Emu(x), y + Inches(0.15), stage_w, Inches(0.4),
             f"STAGE {n}", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tbox(s, Emu(x), y + Inches(0.55), stage_w, Inches(0.4),
             title, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tbox(s, Emu(x) + Inches(0.1), y + Inches(1.05), stage_w - Inches(0.2), Inches(1.0),
             body, size=10, color=WHITE, align=PP_ALIGN.CENTER)
        # Arrow between stages
        if i < 4:
            arrow_x = Emu(x) + stage_w
            tbox(s, arrow_x, y + Inches(0.85), Inches(0.15), Inches(0.4),
                 "▶", size=14, color=NAVY, align=PP_ALIGN.CENTER)

    # Outputs row
    rrect(s, Inches(0.6), Inches(3.7), prs.slide_width - Inches(1.2), Inches(1.4), PALE)
    tbox(s, Inches(0.85), Inches(3.85), Inches(6), Inches(0.35),
         "OUTPUTS", size=11, bold=True, color=CYAN)
    multi_tbox(
        s, Inches(0.85), Inches(4.15), Inches(12), Inches(0.95),
        [
            ("Per-app JSON:", "raw evidence — model artefacts, entropy, permissions, Smali hits, coverage stats"),
            ("Aggregated CSV:", "one row per app — severity counts, top finding, SHA-256 APK hash"),
            ("Severity rubric:", "Critical / High / Medium / Info — assigned per finding, propagated to app-level"),
        ],
        size=11, color=GREY,
    )

    # Disclosure box
    rrect(s, Inches(0.6), Inches(5.3), prs.slide_width - Inches(1.2), Inches(1.5), NAVY)
    tbox(s, Inches(0.85), Inches(5.4), Inches(12), Inches(0.3),
         "HONEST SCOPE DISCLOSURE", size=11, bold=True, color=AMBER)
    tbox(s, Inches(0.85), Inches(5.7), Inches(12), Inches(0.4),
         "Stage 5 is pattern-matching (regex on Smali), not taint analysis.",
         size=14, bold=True, color=WHITE)
    tbox(s, Inches(0.85), Inches(6.1), Inches(12), Inches(0.65),
         "Full inter-procedural Soot/FlowDroid taint was planned but deferred. ProGuard / R8 obfuscation "
         "can evade literal-string patterns — this is the dominant source of false-negative risk and is "
         "disclosed in limitations.",
         size=11, color=LIGHT_BLUE, italic=True)

    page_num(s, prs, 5)


def slide_6_data(prs):
    s = blank(prs)
    header(s, prs, "Methods: Data & Severity Rubric",
           "Pilot corpus of six APKs — five F-Droid open-source apps plus a synthetic positive control")

    # Left: corpus table
    tbox(s, Inches(0.6), Inches(1.1), Inches(8), Inches(0.35),
         "PILOT CORPUS", size=11, bold=True, color=CYAN)
    tbox(s, Inches(0.6), Inches(1.4), Inches(8), Inches(0.35),
         "F-Droid chosen for license friction — no Play Store ToS issues. Synthetic positive control = "
         "stock MobileNet v1 in an APK-shaped zip.",
         size=10, color=GREY, italic=True)

    # Table
    col_widths = [Inches(2.4), Inches(1.0), Inches(1.8), Inches(2.0)]
    header_row_y = Inches(2.05)
    headers = ["App", "Size", "SHA-256 (prefix)", "Role"]
    # Header
    rect(s, Inches(0.6), header_row_y, sum(col_widths, Emu(0)), Inches(0.4), NAVY)
    x = Inches(0.6)
    for i, h in enumerate(headers):
        tbox(s, x + Inches(0.1), header_row_y + Inches(0.08),
             col_widths[i] - Inches(0.2), Inches(0.28),
             h, size=11, bold=True, color=WHITE)
        x += col_widths[i]

    corpus = [
        ("F-Droid (org.fdroid.fdroid)", "12.4 MB", "985f5181d48bb6ba", "App catalogue"),
        ("NewPipe (org.schabi.newpipe)", "10.9 MB", "dbc8a1bb7a3db16f", "YouTube frontend"),
        ("Aves Gallery libre", "55.9 MB", "75697b19f2eb850f", "Photo viewer"),
        ("Nextcloud Files", "77.8 MB", "e36e6ef4215cf003", "Cloud client"),
        ("Fennec (Firefox fork)", "117.6 MB", "251146a2b5f6d801", "Web browser"),
        ("positive_control (MobileNet v1)", "15.7 MB", "d9c71359fcce2d5f", "★ synthetic control"),
    ]
    for i, row in enumerate(corpus):
        row_y = header_row_y + Inches(0.4) + Inches(0.4) * i
        fill = PALE if i % 2 == 0 else WHITE
        rect(s, Inches(0.6), row_y, sum(col_widths, Emu(0)), Inches(0.4), fill)
        x = Inches(0.6)
        for j, val in enumerate(row):
            is_control = i == 5
            color = CYAN if is_control and j == 0 else NAVY
            bold = is_control and j == 0
            tbox(s, x + Inches(0.1), row_y + Inches(0.08),
                 col_widths[j] - Inches(0.2), Inches(0.28),
                 val, size=10, bold=bold, color=color)
            x += col_widths[j]

    # Right: severity rubric
    tbox(s, Inches(9.2), Inches(1.1), Inches(3.6), Inches(0.35),
         "SEVERITY RUBRIC (as implemented)", size=11, bold=True, color=CYAN)
    tbox(s, Inches(9.2), Inches(1.4), Inches(3.6), Inches(0.35),
         "Pilot produced only High + Info (see notes)",
         size=10, color=GREY, italic=True)
    rubric = [
        ("Critical", "Active poisoning surface + unauth model update",
         "Not producible by pilot (RQ4 deferred)", RED, True),
        ("High", "Identified unencrypted model OR poisoning-surface regex hit",
         "Fires on positive control", ORANGE, False),
        ("Medium", "Model file with ambiguous entropy (encryption vs. compression)",
         "Not produced by pilot corpus", AMBER, True),
        ("Info", "No on-device ML detected — outside Truth-Serum surface",
         "Fires on all F-Droid apps", GREY, False),
    ]
    y = Inches(1.95)
    for label, desc, note, color, theoretical in rubric:
        rrect(s, Inches(9.2), y, Inches(3.6), Inches(1.05), WHITE)
        border = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(9.2), y, Inches(3.6), Inches(1.05))
        border.fill.background()
        border.line.color.rgb = LIGHT_GREY
        border.line.width = Pt(1)
        # Colored strip
        rect(s, Inches(9.2), y, Inches(0.2), Inches(1.05), color)
        tbox(s, Inches(9.5), y + Inches(0.08), Inches(3.0), Inches(0.3),
             label, size=13, bold=True, color=color)
        tbox(s, Inches(9.5), y + Inches(0.36), Inches(3.3), Inches(0.4),
             desc, size=9, color=GREY)
        tbox(s, Inches(9.5), y + Inches(0.74), Inches(3.3), Inches(0.28),
             note, size=8, italic=True,
             color=RED if theoretical else GREEN)
        y += Inches(1.13)

    page_num(s, prs, 6)


def slide_7_evaluation(prs):
    s = blank(prs)
    header(s, prs, "Methods: Evaluation Design",
           "Two calibration moves — positive control proves detection fires; null on F-Droid proves no spurious alarms")

    # Top: two panels (positive control + null hypothesis)
    panel_w = Inches(6.05)
    panel_h = Inches(2.9)
    y = Inches(1.25)

    # Positive control panel
    rrect(s, Inches(0.6), y, panel_w, panel_h, WHITE)
    border = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.6), y, panel_w, panel_h)
    border.fill.background()
    border.line.color.rgb = LIGHT_GREY
    border.line.width = Pt(1)
    rect(s, Inches(0.6), y, panel_w, Inches(0.12), GREEN)
    tbox(s, Inches(0.8), y + Inches(0.22), panel_w - Inches(0.4), Inches(0.4),
         "POSITIVE CONTROL", size=12, bold=True, color=GREEN)
    tbox(s, Inches(0.8), y + Inches(0.6), panel_w - Inches(0.4), Inches(0.4),
         "Synthetic APK · MobileNet v1 · 16.9 MB",
         size=14, bold=True, color=NAVY)
    multi_tbox(
        s, Inches(0.8), y + Inches(1.05), panel_w - Inches(0.4), Inches(1.8),
        [
            ("Expected:", "detection fires with High severity"),
            ("Validates:", "Stage 1 (archive scan) — file-extension match, magic bytes, entropy"),
            ("Known weakness:", "zip-with-model, not a valid APK structure — apktool decode fails by design, so Stages 2–5 are untested by this control"),
            ("Future work:", "minimal valid APK + ProGuard-obfuscated variant to exercise Stages 4–5"),
        ],
        size=10, color=GREY,
    )

    # Null hypothesis panel
    rrect(s, Inches(6.9), y, panel_w, panel_h, WHITE)
    border = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(6.9), y, panel_w, panel_h)
    border.fill.background()
    border.line.color.rgb = LIGHT_GREY
    border.line.width = Pt(1)
    rect(s, Inches(6.9), y, panel_w, Inches(0.12), CYAN)
    tbox(s, Inches(7.1), y + Inches(0.22), panel_w - Inches(0.4), Inches(0.4),
         "NULL HYPOTHESIS (F-DROID)", size=12, bold=True, color=CYAN)
    tbox(s, Inches(7.1), y + Inches(0.6), panel_w - Inches(0.4), Inches(0.4),
         "5 open-source utility apps · 274 MB",
         size=14, bold=True, color=NAVY)
    multi_tbox(
        s, Inches(7.1), y + Inches(1.05), panel_w - Inches(0.4), Inches(1.8),
        [
            ("Expected:", "no on-device ML in unobfuscated reproducible builds"),
            ("Validates:", "no false-positive at low baseline — detector correctly outputs null"),
            ("Does NOT validate:", "performance on obfuscated Play Store targets"),
            ("Sample note:", "n=5 ≈ 0.1% of F-Droid catalogue — pilot, not structural claim"),
        ],
        size=10, color=GREY,
    )

    # Bottom: what the two-arm design does / doesn't validate
    rrect(s, Inches(0.6), Inches(4.4), prs.slide_width - Inches(1.2), Inches(2.3), NAVY)
    tbox(s, Inches(0.85), Inches(4.55), Inches(12), Inches(0.3),
         "WHAT THIS EVAL DESIGN ACTUALLY TESTS", size=11, bold=True, color=AMBER)

    left_x = Inches(0.85)
    right_x = Inches(7.0)
    yy = Inches(4.95)
    tbox(s, left_x, yy, Inches(5.8), Inches(0.35),
         "✓ CONFIRMED", size=12, bold=True, color=GREEN)
    multi_tbox(
        s, left_x, yy + Inches(0.32), Inches(5.8), Inches(1.5),
        [
            "Archive stage fires on known model artefact",
            "Severity rubric assigns High to exposed model",
            "Null on clean corpus → no false-positive at baseline",
            "Pipeline is re-runnable end-to-end on new APKs",
        ],
        size=11, color=WHITE,
    )

    tbox(s, right_x, yy, Inches(5.8), Inches(0.35),
         "✗ NOT YET VALIDATED", size=12, bold=True, color=AMBER)
    multi_tbox(
        s, right_x, yy + Inches(0.32), Inches(5.8), Inches(1.5),
        [
            "Poisoning-surface regexes have never fired true-positive",
            "Obfuscation-resistance (ProGuard / R8 renaming)",
            "False-positive / false-negative rates (n=5 too small)",
            "Performance vs. MobSF / ModelXRay (no baseline run)",
        ],
        size=11, color=WHITE,
    )

    page_num(s, prs, 7)


def slide_8_results(prs, rows, agg):
    s = blank(prs)
    header(s, prs, "Results: Findings",
           "Six APKs processed · positive control detected at High severity · null on five F-Droid apps")

    # KPI strip
    y = Inches(1.15)
    kpi_w = Inches(2.9)
    kpi_h = Inches(1.15)
    kpis = [
        (f"{agg['n_apks']}", "APKs analysed", CYAN),
        (f"{agg['total_mb']:.0f} MB", "Total corpus", TEAL),
        (f"{agg['n_with_models']} / {agg['n_apks']}", "With on-device ML", AMBER),
        (f"{agg['total_findings']}", "Findings logged", ORANGE),
    ]
    for i, (num, lbl, color) in enumerate(kpis):
        x = Inches(0.6) + (kpi_w + Inches(0.12)) * i
        kpi(s, x, y, kpi_w, kpi_h, num, lbl, color)

    # Results table
    ty = Inches(2.55)
    tbox(s, Inches(0.6), ty, Inches(12), Inches(0.35),
         "PER-APP FINDINGS", size=11, bold=True, color=CYAN)

    col_widths = [Inches(3.8), Inches(1.0), Inches(1.2), Inches(1.4), Inches(1.6), Inches(1.3), Inches(1.2)]
    header_row_y = ty + Inches(0.4)
    headers = ["App (package)", "Size", "Models", "ML SDKs", "Smali coverage", "Findings", "Top severity"]

    # Header row
    rect(s, Inches(0.6), header_row_y, sum(col_widths, Emu(0)), Inches(0.4), NAVY)
    x = Inches(0.6)
    for i, h in enumerate(headers):
        tbox(s, x + Inches(0.08), header_row_y + Inches(0.09),
             col_widths[i] - Inches(0.16), Inches(0.28),
             h, size=10, bold=True, color=WHITE)
        x += col_widths[i]

    sev_color = {"Critical": RED, "High": ORANGE, "Medium": AMBER, "Info": GREY}

    for i, row in enumerate(rows):
        row_y = header_row_y + Inches(0.4) + Inches(0.36) * i
        fill = PALE if i % 2 == 0 else WHITE
        rect(s, Inches(0.6), row_y, sum(col_widths, Emu(0)), Inches(0.36), fill)

        top_sev = "Info"
        for key in ("critical", "high", "medium", "info"):
            if int(row.get(key, 0)) > 0:
                top_sev = key.capitalize()
                break

        pkg = row.get("package") or Path(row["apk"]).stem
        scanned = int(row.get("smali_scanned", 0) or 0)
        total = int(row.get("smali_total", 0) or 0)
        cov = "n/a" if total == 0 else f"{scanned:,}/{total:,} ({scanned*100//total}%)"
        raw_sdks = (row.get("ml_sdks") or "").strip()
        sdks = raw_sdks if raw_sdks else ("n/a" if scanned == 0 else "none")
        total_findings = sum(int(row.get(k, 0)) for k in ("critical", "high", "medium", "info"))

        cells = [
            pkg,
            f"{row.get('apk_size_mb', '0')} MB",
            row.get("n_model_artifacts", "0"),
            sdks,
            cov,
            str(total_findings),
            top_sev,
        ]
        x = Inches(0.6)
        for j, val in enumerate(cells):
            color = sev_color[top_sev] if j == len(cells) - 1 else NAVY
            bold = j == len(cells) - 1
            tbox(s, x + Inches(0.08), row_y + Inches(0.07),
                 col_widths[j] - Inches(0.16), Inches(0.26),
                 val, size=10, bold=bold, color=color)
            x += col_widths[j]

    # Takeaway banner
    banner_y = Inches(5.7)
    rrect(s, Inches(0.6), banner_y, prs.slide_width - Inches(1.2), Inches(1.1), NAVY)
    tbox(s, Inches(0.85), banner_y + Inches(0.1), Inches(12), Inches(0.3),
         "HEADLINE", size=11, bold=True, color=CYAN)
    tbox(s, Inches(0.85), banner_y + Inches(0.38), Inches(12), Inches(0.45),
         "Detector fires on positive control (High); null on all five F-Droid apps.",
         size=14, bold=True, color=WHITE)
    tbox(s, Inches(0.85), banner_y + Inches(0.75), Inches(12), Inches(0.35),
         "Pilot observation: unobfuscated open-source utility apps don't bundle on-device ML. 0.1% sample — not a structural claim about F-Droid.",
         size=11, color=LIGHT_BLUE, italic=True)

    page_num(s, prs, 8)


def slide_9_coverage_chart(prs, rows):
    s = blank(prs)
    header(s, prs, "Results: Smali Coverage — the Dominant FN Risk",
           "5,000-file cap meant large apps were only partially scanned. Raising it is the single highest-impact engineering change.")

    # Filter to apps with apktool success
    scanned_rows = [r for r in rows if int(r.get("smali_total", 0) or 0) > 0]
    scanned_rows.sort(key=lambda r: int(r["smali_total"]), reverse=True)

    # Build native pptx bar chart
    chart_data = CategoryChartData()
    labels = [Path(r["apk"]).stem.replace("_", " ") for r in scanned_rows]
    coverage_pct = [
        round(int(r["smali_scanned"]) * 100.0 / int(r["smali_total"]), 1)
        for r in scanned_rows
    ]
    chart_data.categories = labels
    chart_data.add_series("Smali coverage (%)", coverage_pct)

    chart_left = Inches(0.6)
    chart_top = Inches(1.2)
    chart_w = Inches(7.5)
    chart_h = Inches(4.6)
    chart_shape = s.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, chart_left, chart_top, chart_w, chart_h, chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    chart.has_title = True
    chart.chart_title.text_frame.text = "Smali scan coverage per app (% of files scanned under 5,000-file cap)"
    for para in chart.chart_title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = NAVY

    # Colour bars by coverage threshold
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.size = Pt(10)
    dl.font.bold = True
    dl.font.color.rgb = NAVY
    dl.number_format = '0.0"%"'

    series = plot.series[0]
    for idx, pct in enumerate(coverage_pct):
        pt = series.points[idx]
        pt.format.fill.solid()
        if pct < 25:
            pt.format.fill.fore_color.rgb = RED
        elif pct < 60:
            pt.format.fill.fore_color.rgb = ORANGE
        else:
            pt.format.fill.fore_color.rgb = GREEN
        pt.format.line.fill.background()

    # Axes styling
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(10)
    cat_axis.tick_labels.font.color.rgb = NAVY
    val_axis = chart.value_axis
    val_axis.maximum_scale = 100
    val_axis.minimum_scale = 0
    val_axis.tick_labels.font.size = Pt(9)
    val_axis.tick_labels.font.color.rgb = GREY

    # Right panel: interpretation
    legend_x = Inches(8.4)
    legend_y = Inches(1.25)
    tbox(s, legend_x, legend_y, Inches(4.5), Inches(0.35),
         "INTERPRETATION", size=11, bold=True, color=CYAN)

    legend_items = [
        ("< 25%", RED, "Large apps — high false-negative risk. Ninefold increase would reach full coverage."),
        ("25–60%", ORANGE, "Moderate apps — residual FN risk. Priority-ordered scanning would help."),
        ("> 60%", GREEN, "Small apps — near-full coverage. Null is meaningful here."),
    ]
    y = Inches(1.6)
    for label, color, desc in legend_items:
        rrect(s, legend_x, y, Inches(4.5), Inches(0.95), WHITE)
        border = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    legend_x, y, Inches(4.5), Inches(0.95))
        border.fill.background()
        border.line.color.rgb = LIGHT_GREY
        border.line.width = Pt(1)
        rect(s, legend_x, y, Inches(0.25), Inches(0.95), color)
        tbox(s, legend_x + Inches(0.4), y + Inches(0.1), Inches(4.0), Inches(0.3),
             label, size=12, bold=True, color=color)
        tbox(s, legend_x + Inches(0.4), y + Inches(0.38), Inches(4.0), Inches(0.5),
             desc, size=10, color=GREY)
        y += Inches(1.05)

    # Bottom takeaway
    rrect(s, Inches(0.6), Inches(6.05), prs.slide_width - Inches(1.2), Inches(0.95), NAVY)
    tbox(s, Inches(0.85), Inches(6.15), Inches(12), Inches(0.3),
         "ACTIONABLE FIX", size=11, bold=True, color=AMBER)
    tbox(s, Inches(0.85), Inches(6.45), Inches(12), Inches(0.5),
         "Raise the Smali cap OR switch to priority-ordered scanning (SDK-import-first). "
         "Trivial code change; not in tonight's build. Materially reduces FN risk on large apps.",
         size=11, color=WHITE)

    page_num(s, prs, 9)


def slide_10_limitations_conclusion(prs):
    s = blank(prs)
    header(s, prs, "Limitations & Conclusion",
           "What this pilot can and cannot claim, and where the tool goes next")

    # Limitations (left) — tighter
    tbox(s, Inches(0.6), Inches(1.1), Inches(6.2), Inches(0.35),
         "LIMITATIONS DISCLOSED UP-FRONT", size=11, bold=True, color=RED)
    multi_tbox(
        s, Inches(0.6), Inches(1.45), Inches(6.2), Inches(3.0),
        [
            ("Pattern-match, not taint analysis.", "Regex on Smali — fragile against ProGuard / R8 class renaming."),
            ("n = 5 corpus.", "Pilot scale. Cannot compute FP / FN rates from this sample."),
            ("Positive control validates only Stage 1.", "Archive match on unobfuscated model. SDK / Smali / severity untested."),
            ("No baseline comparison.", "MobSF, QARK, ModelXRay not run on same corpus."),
            ("No manual ground truth.", "Null results not confirmed by jadx inspection."),
            ("Smali cap at 5,000 files.", "Large apps 11–25% covered — see slide 9."),
        ],
        size=10, color=GREY,
    )

    # Next steps (right) — tighter
    tbox(s, Inches(7.0), Inches(1.1), Inches(5.8), Inches(0.35),
         "CONCRETE NEXT STEPS", size=11, bold=True, color=GREEN)
    multi_tbox(
        s, Inches(7.0), Inches(1.45), Inches(5.8), Inches(3.0),
        [
            ("Play-Store corpus (n≈100).", "Keyboards, translation, camera — where training surface lives."),
            ("Obfuscation-resistant detector.", "Anchor on native-lib presence + string-constant extraction."),
            ("Baseline run.", "MobSF + ModelXRay on same apps → comparative FP/FN rates."),
            ("Manual verification loop.", "Spot-check 10–20 apps in jadx to anchor null results."),
            ("Raise Smali cap.", "Direct FN-risk reduction on large apps (slide 9)."),
            ("Stronger positive control.", "Valid minimal APK + ProGuard-obfuscated variant."),
        ],
        size=10, color=GREY,
    )

    # TAKE-AWAY banner — the single quotable sentence
    rrect(s, Inches(0.6), Inches(4.55), prs.slide_width - Inches(1.2), Inches(0.7), NAVY)
    rect(s, Inches(0.6), Inches(4.55), Inches(0.15), Inches(0.7), AMBER)
    tbox(s, Inches(0.9), Inches(4.62), Inches(2.3), Inches(0.28),
         "TAKE-AWAY", size=10, bold=True, color=AMBER)
    tbox(s, Inches(0.9), Inches(4.87), prs.slide_width - Inches(1.5), Inches(0.35),
         "The tool is calibrated. The science starts at the Play Store corpus.",
         size=16, bold=True, color=WHITE)

    # Contribution banner (compact)
    rrect(s, Inches(0.6), Inches(5.35), prs.slide_width - Inches(1.2), Inches(1.5), NAVY)
    tbox(s, Inches(0.85), Inches(5.45), Inches(12), Inches(0.3),
         "CONTRIBUTIONS", size=11, bold=True, color=CYAN)

    contribs = [
        ("1", "Static-analysis pipeline targeting the training-surface gap orthogonal to Xu '19 and Sun '21.", CYAN),
        ("2", "Severity-rubric design tied to the Truth-Serum precondition — High and Info demonstrated; Critical and Medium reserved for future detectors.", TEAL),
        ("3", "Calibrated pilot: positive control fires as predicted; null on 5 open-source apps with documented coverage gaps.", AMBER),
    ]
    cont_w = Inches(4.03)
    cont_h = Inches(1.0)
    cy = Inches(5.78)
    for i, (n, text, color) in enumerate(contribs):
        x = Inches(0.85) + (cont_w + Inches(0.08)) * i
        rrect(s, x, cy, cont_w, cont_h, color)
        tbox(s, x, cy + Inches(0.05), cont_w, Inches(0.3),
             n, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tbox(s, x + Inches(0.12), cy + Inches(0.38), cont_w - Inches(0.24), cont_h - Inches(0.44),
             text, size=9, color=WHITE, align=PP_ALIGN.CENTER)

    # Reproducibility footer
    tbox(s, Inches(0.6), Inches(6.95), prs.slide_width - Inches(1.2), Inches(0.22),
         "REPRODUCIBILITY   ·   python3 pipeline/pipeline.py --apk-dir apks/ --out results/   ·   APK SHA-256 prefixes on slide 6   ·   source + JSON/CSV evidence archived with submission",
         size=9, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    page_num(s, prs, 10)


def slide_11_qa(prs):
    """Appendix — anticipated questions with pre-rehearsed answers."""
    s = blank(prs)
    header(s, prs, "Appendix: Anticipated Q & A",
           "Preempted answers — shown only if questions arise. Rehearsed talk does not walk through this slide.")

    qa = [
        ("Would your tool detect federated learning in GBoard?",
         "Not reliably. ProGuard / R8 renames the class names my regexes target. The tool would flag the RECORD_AUDIO permission + libtensorflowlite.so, but not the federated training surface itself. Obfuscation-resistant detection (anchored on native-lib presence + string-constant extraction) is a concrete next step."),
        ("How is this different from ModelXRay (Sun '21)?",
         "ModelXRay audits model protection — can the model be stolen? This audits the training surface — is user input fed into training? Orthogonal threat dimensions. On this pilot corpus, the one detection that fired (positive control, model exposure) is in fact within ModelXRay's scope; the differentiation is architectural and demonstrates on future Play Store corpora."),
        ("Is your synthetic positive control really a valid APK?",
         "No — it is a zip containing a MobileNet .tflite file, not an APK with a manifest / dex / Smali. That is why apktool decode fails and why stages 2–5 show n/a on slide 8. The positive control validates stage 1 only. A valid minimal APK + a ProGuard-obfuscated variant are on the future-work list."),
        ("Your rubric defines 4 levels but produces only 2. Why not a 2-level rubric?",
         "Critical and Medium anticipate detectors not yet built — insecure model-update channels (RQ4, deferred) and partial-encryption entropy signals. The 4-level design is intentional. The pilot demonstrates 2 of 4 in practice; slide 6 discloses this explicitly."),
        ("On this 5-app corpus, what does your tool do that `grep -r tflite` cannot?",
         "On this corpus alone — materially very little; the archive-stage detection is the only finding. The tool's differentiation (rubric, threat-model-derived severity, stage-5 pattern scan) is architectural, and is validated on the positive control's archive detection. Empirical differentiation on Play Store apps is future work."),
        ("FP / FN rates?",
         "Undefined with n = 5. A labelled ground-truth corpus of ≈100 apps — manually verified in jadx — is required for rate estimation. Listed as future work on slide 10."),
    ]

    y = Inches(1.15)
    row_h = Inches(0.87)
    for q, a in qa:
        rrect(s, Inches(0.6), y, prs.slide_width - Inches(1.2), row_h, WHITE)
        border = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.6), y,
                                    prs.slide_width - Inches(1.2), row_h)
        border.fill.background()
        border.line.color.rgb = LIGHT_GREY
        border.line.width = Pt(1)
        # Accent stripe (left)
        rect(s, Inches(0.6), y, Inches(0.12), row_h, CYAN)
        # Question
        tbox(s, Inches(0.85), y + Inches(0.07),
             prs.slide_width - Inches(1.65), Inches(0.3),
             "Q.  " + q, size=11, bold=True, color=NAVY)
        # Answer
        tbox(s, Inches(0.85), y + Inches(0.37),
             prs.slide_width - Inches(1.65), row_h - Inches(0.42),
             "A.  " + a, size=9, color=GREY)
        y += row_h + Inches(0.04)

    page_num(s, prs, 11)


# ----------------------------------------------------------------------

def main():
    prs = new_prs()
    rows = load_rows()
    agg = aggregate(rows)

    slide_1_title(prs)
    slide_2_motivation(prs)
    slide_3_related_work(prs)
    slide_4_threat_model(prs)
    slide_5_approach(prs)
    slide_6_data(prs)
    slide_7_evaluation(prs)
    slide_8_results(prs, rows, agg)
    slide_9_coverage_chart(prs, rows)
    slide_10_limitations_conclusion(prs)
    slide_11_qa(prs)

    DST.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(DST))
    print(f"saved: {DST}")
    print(f"   {len(prs.slides)} slides total")


if __name__ == "__main__":
    main()
